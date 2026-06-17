"""DeckSpec → python-pptx Presentation 的适配器。

把 ljg-ppt-design 出的 JSON 渲染成真实 .pptx 文件,跨平台 (Mac/Windows/Linux)。
依赖: python-pptx (可选扩展)。

用法:
  from ljg_ppt_design import render_deck
  from ljg_ppt_design.data.pptx_renderer import render_to_pptx
  deck = render_deck("academic", "school", content)
  render_to_pptx(deck, "/tmp/output.pptx")
"""

from __future__ import annotations

from typing import Any

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.oxml.ns import qn
except ImportError as e:
    raise ImportError("pptx_renderer 需要 python-pptx,装一下: pip install python-pptx") from e

from ..design_presets import DesignPreset, get_preset
from ..slide_spec import DeckSpec, SlideSpec, SlideElement


EMU_PER_INCH = 914400
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
CANVAS_PX_W = 960
CANVAS_PX_H = 540


def px_to_in_x(px: int) -> float:
    return px / CANVAS_PX_W * SLIDE_W_IN


def px_to_in_y(px: int) -> float:
    return px / CANVAS_PX_H * SLIDE_H_IN


def to_rgb(preset: DesignPreset, role: str) -> RGBColor:
    return RGBColor(*preset.get_color(role))


def font_size_for(role: str, layout_default: int = 18) -> int:
    return {"title": 40, "subtitle": 24, "body": 20, "caption": 14}.get(role, layout_default)


def _set_slide_bg(slide, color: RGBColor) -> None:
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(slide, el: SlideElement, preset: DesignPreset) -> None:
    if not el.text:
        return
    tx = slide.shapes.add_textbox(
        Inches(px_to_in_x(el.x)), Inches(px_to_in_y(el.y)),
        Inches(px_to_in_x(el.w)), Inches(px_to_in_y(el.h)),
    )
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)

    p = tf.paragraphs[0]
    p.alignment = {0: PP_ALIGN.LEFT, 1: PP_ALIGN.CENTER, 2: PP_ALIGN.RIGHT}.get(el.align, PP_ALIGN.LEFT)

    run = p.add_run()
    run.text = el.text
    run.font.size = Pt(el.fs if el.fs > 0 else font_size_for(el.role))
    run.font.bold = el.bold
    run.font.color.rgb = to_rgb(preset, el.color)
    if any('一' <= c <= '鿿' for c in el.text):
        run.font.name = "Microsoft YaHei"
        rPr = run._r.get_or_add_rPr()
        from lxml import etree
        ea = rPr.find(qn("a:ea"))
        if ea is None:
            ea = etree.SubElement(rPr, qn("a:ea"))
        ea.set("typeface", "Microsoft YaHei")


def _add_line(slide, el: SlideElement, preset: DesignPreset) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(px_to_in_x(el.x)), Inches(px_to_in_y(el.y)),
        Inches(px_to_in_x(max(el.w, 1))), Inches(px_to_in_y(max(el.h, 1))),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = to_rgb(preset, el.color)
    shape.line.fill.background()


def _add_box(slide, el: SlideElement, preset: DesignPreset) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(px_to_in_x(el.x)), Inches(px_to_in_y(el.y)),
        Inches(px_to_in_x(el.w)), Inches(px_to_in_y(el.h)),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = to_rgb(preset, el.color)
    shape.line.fill.background()
    if el.text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = el.text
        run.font.size = Pt(16)
        run.font.color.rgb = to_rgb(preset, "dark")


def _add_image_placeholder(slide, el: SlideElement, preset: DesignPreset) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(px_to_in_x(el.x)), Inches(px_to_in_y(el.y)),
        Inches(px_to_in_x(el.w)), Inches(px_to_in_y(el.h)),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = to_rgb(preset, "light")
    shape.line.color.rgb = to_rgb(preset, "gray")
    shape.line.width = Pt(1)
    if el.text:
        tf = shape.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f"[ {el.text} ]"
        run.font.size = Pt(14)
        run.font.italic = True
        run.font.color.rgb = to_rgb(preset, "gray")


def _add_grid_cards(slide, el: SlideElement, items: list, preset: DesignPreset) -> None:
    if not items:
        items = [{}]
    cols = max(el.cols, 1)
    rows = max(el.rows, 1) if el.rows else 1
    item_w = el.item_w or (el.w // cols)
    item_h = el.item_h or (el.h // rows)

    for idx, item in enumerate(items):
        c = idx % cols
        r = idx // cols
        if r >= rows:
            break
        ix = el.x + c * (item_w + el.gap_x)
        iy = el.y + r * (item_h + el.gap_y)
        if item_w > 0 and item_h > 0:
            card_bg = el.styles[0].get("bg") if el.styles else None
            if card_bg and card_bg.startswith("{"):
                color_key = card_bg.strip("{}")
                card_bg = item.get(color_key, "primary")
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(px_to_in_x(ix)), Inches(px_to_in_y(iy)),
                Inches(px_to_in_x(item_w)), Inches(px_to_in_y(item_h)),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = to_rgb(preset, card_bg or "light")
            shape.line.fill.background()
        if el.fields:
            for fi, field_tpl in enumerate(el.fields):
                field_key = field_tpl.strip("{}")
                value = str(item.get(field_key, ""))
                style = el.styles[fi] if fi < len(el.styles) else {}
                fs = style.get("fs", 16)
                bold = style.get("bold", False)
                color_role = style.get("color", "dark")
                align = {0: PP_ALIGN.LEFT, 1: PP_ALIGN.CENTER, 2: PP_ALIGN.RIGHT}.get(
                    style.get("align", 0), PP_ALIGN.LEFT)
                y_offset = style.get("y_offset", 0)
                w = style.get("w", item_w)
                h = style.get("h", item_h - y_offset)
                if style.get("bg") and not style["bg"].startswith("{"):
                    bar = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(px_to_in_x(ix)), Inches(px_to_in_y(iy + y_offset)),
                        Inches(px_to_in_x(w)), Inches(px_to_in_y(h)),
                    )
                    bar.fill.solid()
                    bar.fill.fore_color.rgb = to_rgb(preset, style["bg"])
                    bar.line.fill.background()
                tx = slide.shapes.add_textbox(
                    Inches(px_to_in_x(ix + 4)), Inches(px_to_in_y(iy + y_offset + 4)),
                    Inches(px_to_in_x(w - 8)), Inches(px_to_in_y(h - 8)),
                )
                tf = tx.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.alignment = align
                run = p.add_run()
                run.text = value
                run.font.size = Pt(fs)
                run.font.bold = bold
                run.font.color.rgb = to_rgb(preset, color_role)


def _add_timeline_items(slide, el: SlideElement, items: list, preset: DesignPreset) -> None:
    if not items:
        return
    count = min(len(items), el.count or len(items))
    for i, item in enumerate(items[:count]):
        iy = el.start_y + i * el.gap
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(px_to_in_x(el.x - 8)), Inches(px_to_in_y(iy + 4)),
            Inches(px_to_in_x(12)), Inches(px_to_in_y(12)),
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = to_rgb(preset, "primary")
        dot.line.fill.background()
        if el.styles and len(el.styles) >= 1:
            s = el.styles[0]
            tx = slide.shapes.add_textbox(
                Inches(px_to_in_x(el.x + 10)), Inches(px_to_in_y(iy)),
                Inches(px_to_in_x(s.get("w", 110))), Inches(px_to_in_y(30)),
            )
            p = tx.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(item.get("date", ""))
            run.font.size = Pt(s.get("fs", 16))
            run.font.bold = s.get("bold", True)
            run.font.color.rgb = to_rgb(preset, s.get("color", "primary"))
        if el.styles and len(el.styles) >= 2:
            s = el.styles[1]
            tx2 = slide.shapes.add_textbox(
                Inches(px_to_in_x(el.x + 130)), Inches(px_to_in_y(iy)),
                Inches(px_to_in_x(s.get("w", 550))), Inches(px_to_in_y(40)),
            )
            p = tx2.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(item.get("event", ""))
            run.font.size = Pt(s.get("fs", 16))
            run.font.color.rgb = to_rgb(preset, s.get("color", "dark"))


def _add_pipeline_items(slide, el: SlideElement, items: list, preset: DesignPreset) -> None:
    if not items:
        return
    for i, item in enumerate(items):
        ix = el.x + i * (el.item_w + el.gap)
        iy = el.y
        color_role = item.get("color", ["primary", "secondary", "accent", "dark"][i % 4])
        block = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(px_to_in_x(ix)), Inches(px_to_in_y(iy)),
            Inches(px_to_in_x(el.item_w)), Inches(px_to_in_y(50)),
        )
        block.fill.solid()
        block.fill.fore_color.rgb = to_rgb(preset, color_role)
        block.line.fill.background()
        tx = slide.shapes.add_textbox(
            Inches(px_to_in_x(ix + 4)), Inches(px_to_in_y(iy + 10)),
            Inches(px_to_in_x(el.item_w - 8)), Inches(px_to_in_y(30)),
        )
        p = tx.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(item.get("step_num", str(i + 1)))
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = to_rgb(preset, "white")
        tx2 = slide.shapes.add_textbox(
            Inches(px_to_in_x(ix + 4)), Inches(px_to_in_y(iy + 60)),
            Inches(px_to_in_x(el.item_w - 8)), Inches(px_to_in_y(30)),
        )
        p = tx2.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(item.get("step_name", ""))
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = to_rgb(preset, "primary")
        if item.get("step_detail"):
            tx3 = slide.shapes.add_textbox(
                Inches(px_to_in_x(ix + 4)), Inches(px_to_in_y(iy + 90)),
                Inches(px_to_in_x(el.item_w - 8)), Inches(px_to_in_y(150)),
            )
            tf = tx3.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = str(item.get("step_detail", ""))
            run.font.size = Pt(10)
            run.font.color.rgb = to_rgb(preset, "dark")


def _add_table(slide, el: SlideElement, data: Any, preset: DesignPreset) -> None:
    if not isinstance(data, dict):
        return
    rows = el.rows or 2
    cols = el.cols or 2
    table_data = data.get("table_data", [])
    actual_rows = min(len(table_data) + 1, rows)
    actual_cols = cols

    shape = slide.shapes.add_table(
        actual_rows, actual_cols,
        Inches(px_to_in_x(el.x)), Inches(px_to_in_y(el.y)),
        Inches(px_to_in_x(el.w - 220)),
        Inches(px_to_in_y(min(actual_rows * (el.row_h or 42), el.h))),
    )
    table = shape.table
    for c in range(actual_cols):
        cell = table.cell(0, c)
        cell.text = str(data.get(f"col{c+1}_name", f"Col {c+1}"))
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(14)
                run.font.bold = True
                run.font.color.rgb = to_rgb(preset, "white")
        cell.fill.solid()
        cell.fill.fore_color.rgb = to_rgb(preset, "primary")
    for r, row_data in enumerate(table_data[:actual_rows - 1], start=1):
        for c in range(min(len(row_data), actual_cols)):
            cell = table.cell(r, c)
            cell.text = str(row_data[c])
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(12)
                    run.font.color.rgb = to_rgb(preset, "dark")


def _render_slide(slide, slide_spec: SlideSpec, preset: DesignPreset) -> None:
    raw_data = slide_spec.data
    is_list_data = isinstance(raw_data, list)
    data = raw_data if isinstance(raw_data, dict) else {}
    items_list = data.get("items") if isinstance(data.get("items"), list) else (
        raw_data if is_list_data else None)

    for el in slide_spec.elements:
        if el.type == "bg":
            _set_slide_bg(slide, to_rgb(preset, el.color))
        elif el.type == "line":
            _add_line(slide, el, preset)
        elif el.type == "text":
            _add_text(slide, el, preset)
        elif el.type == "box":
            _add_box(slide, el, preset)
        elif el.type == "image_placeholder":
            _add_image_placeholder(slide, el, preset)
        elif el.type in ("grid_cards", "grid_items"):
            cards = items_list if items_list is not None else []
            _add_grid_cards(slide, el, cards, preset)
        elif el.type == "timeline_items":
            _add_timeline_items(slide, el, items_list or [], preset)
        elif el.type == "pipeline_items":
            _add_pipeline_items(slide, el, items_list or [], preset)
        elif el.type == "sidebar":
            _add_box(slide, el, preset)
        elif el.type == "table":
            _add_table(slide, el, data, preset)


def render_to_pptx(deck: DeckSpec, output_path: str, preset: str | None = None) -> str:
    preset_key = preset or deck.preset
    p = get_preset(preset_key)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank_layout = prs.slide_layouts[6]

    for slide_spec in deck.slides:
        slide = prs.slides.add_slide(blank_layout)
        _render_slide(slide, slide_spec, p)

    prs.save(output_path)
    return output_path
